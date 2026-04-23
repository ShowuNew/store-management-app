"""
產生「每日店鋪維護系統 - 開發討論會議」簡報
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
import copy

# ── 品牌色 ──────────────────────────────────────────────────────────────────
GREEN      = RGBColor(0x00, 0xA0, 0x40)
DARK_GREEN = RGBColor(0x00, 0x7D, 0x30)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF4, 0xF4, 0xF4)
DARK_GRAY  = RGBColor(0x33, 0x33, 0x33)
MID_GRAY   = RGBColor(0x66, 0x66, 0x66)
ACCENT     = RGBColor(0x00, 0x6B, 0xFF)   # 藍色強調

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── 共用 helper ─────────────────────────────────────────────────────────────
def blank_layout(prs):
    return prs.slide_layouts[6]   # blank

def add_rect(slide, l, t, w, h, fill_color, line_color=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, l, t, w, h, text, font_size=18, bold=False,
                color=DARK_GRAY, align=PP_ALIGN.LEFT, italic=False,
                wrap=True, line_spacing=None):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = wrap
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    if line_spacing:
        from pptx.util import Pt as Ptx
        p.line_spacing = Ptx(line_spacing)
    return tb

def slide_header(slide, title, subtitle=None):
    """頁面頂部綠色橫幅"""
    add_rect(slide, 0, 0, 13.33, 1.15, GREEN)
    add_textbox(slide, 0.4, 0.18, 11, 0.6, title,
                font_size=28, bold=True, color=WHITE)
    if subtitle:
        add_textbox(slide, 0.4, 0.72, 11, 0.38, subtitle,
                    font_size=14, color=RGBColor(0xCC, 0xFF, 0xDD))
    # 右下角頁碼用頁數（在最後統一加）

def slide_footer(slide, page_num, total=9):
    add_textbox(slide, 12.3, 7.1, 1, 0.3,
                f"{page_num} / {total}", font_size=10,
                color=MID_GRAY, align=PP_ALIGN.RIGHT)
    add_textbox(slide, 0.3, 7.1, 4, 0.3,
                "日翊資訊 × FamilyMart 每日店鋪維護系統", font_size=10,
                color=MID_GRAY)

def bullet_lines(slide, l, t, w, items, font_size=16,
                 color=DARK_GRAY, line_h=0.38, dot="●  "):
    for i, item in enumerate(items):
        add_textbox(slide, l, t + i * line_h, w, line_h,
                    dot + item, font_size=font_size, color=color)

# ═══════════════════════════════════════════════════════════════════════════
# 封面
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout(prs))
add_rect(slide, 0, 0, 13.33, 7.5, GREEN)            # 全綠底
add_rect(slide, 0, 4.8, 13.33, 2.7, DARK_GREEN)     # 下深綠

# 白色大標
add_textbox(slide, 1.0, 1.3, 11.3, 1.2,
            "每日店鋪維護系統", font_size=48, bold=True,
            color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide, 1.0, 2.6, 11.3, 0.7,
            "開發討論會議報告", font_size=26,
            color=RGBColor(0xCC, 0xFF, 0xDD), align=PP_ALIGN.CENTER)
add_textbox(slide, 1.0, 3.4, 11.3, 0.55,
            "Daily Store Maintenance System — Development Discussion", font_size=15,
            italic=True, color=RGBColor(0xAA, 0xEE, 0xBB), align=PP_ALIGN.CENTER)

add_textbox(slide, 1.0, 5.3, 11.3, 0.5,
            "日翊資訊  ×  全家便利商店（FamilyMart）", font_size=18,
            color=WHITE, align=PP_ALIGN.CENTER, bold=True)
add_textbox(slide, 1.0, 5.95, 11.3, 0.45,
            "2026 年 4 月 15 日", font_size=15,
            color=RGBColor(0xCC, 0xFF, 0xDD), align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════
# P2 — 議程
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout(prs))
add_rect(slide, 0, 0, 13.33, 7.5, LIGHT_GRAY)
slide_header(slide, "議程", "Agenda")
slide_footer(slide, 2)

items = [
    ("01", "前台操作介面  —  改善項目"),
    ("02", "全管家串接現況"),
    ("03", "角色權限對應總覽"),
    ("04", "IOT 設備串接規劃"),
    ("05", "後台 WEB 報表"),
    ("06", "資料留存時間"),
    ("07", "下一步行動項目"),
]
for idx, (no, text) in enumerate(items):
    row_t = 1.35 + idx * 0.74
    add_rect(slide, 0.6, row_t, 0.72, 0.52, GREEN)
    add_textbox(slide, 0.6, row_t + 0.05, 0.72, 0.44,
                no, font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, 1.5, row_t + 0.08, 10, 0.44,
                text, font_size=17, color=DARK_GRAY)

# ═══════════════════════════════════════════════════════════════════════════
# P3 — 前台操作介面：改善項目
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout(prs))
add_rect(slide, 0, 0, 13.33, 7.5, WHITE)
slide_header(slide, "前台操作介面  —  改善項目", "Frontend UX Improvements")
slide_footer(slide, 3)

cards = [
    ("01", "小店長連結（免登入）",
     "店長可直接產生連結，小店長透過\n連結存取系統，無需帳號密碼登入。"),
    ("02", "設備溫度自動帶入",
     "自動帶入上一筆紀錄溫度；\n若無紀錄則帶入標準預設溫度值。"),
    ("03", "咖啡機自檢項目",
     "新增咖啡機自主檢查表單，\n確保每日機台清潔與品質標準。"),
    ("04", "GPS 定位功能",
     "記錄填寫位置，確認店員於店內\n完成作業，提高紀錄可信度。"),
    ("05", "擔當轄區完成狀況",
     "擔當可即時查看轄區所有店鋪\n每日工作的完成進度概覽。"),
]

col_w = 2.4
col_gap = 0.15
start_l = 0.35
for i, (no, title, desc) in enumerate(cards):
    l = start_l + i * (col_w + col_gap)
    t = 1.3
    # 卡片底
    add_rect(slide, l, t, col_w, 5.6, RGBColor(0xF0, 0xF8, 0xF2),
             line_color=RGBColor(0xCC, 0xEE, 0xD4))
    # 綠色頂條
    add_rect(slide, l, t, col_w, 0.45, GREEN)
    add_textbox(slide, l, t + 0.03, col_w, 0.4,
                no, font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, l + 0.1, t + 0.55, col_w - 0.2, 0.65,
                title, font_size=13, bold=True, color=DARK_GREEN)
    add_textbox(slide, l + 0.1, t + 1.25, col_w - 0.2, 4.1,
                desc, font_size=12, color=DARK_GRAY)

# ═══════════════════════════════════════════════════════════════════════════
# P4 — 全管家串接現況
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout(prs))
add_rect(slide, 0, 0, 13.33, 7.5, WHITE)
slide_header(slide, "全管家串接現況", "FamilyMart Backend Integration Status")
slide_footer(slide, 4)

# 狀態卡 — 進行中
add_rect(slide, 0.5, 1.4, 5.5, 2.8, RGBColor(0xFF, 0xF8, 0xE5),
         line_color=RGBColor(0xFF, 0xCC, 0x00))
add_rect(slide, 0.5, 1.4, 5.5, 0.5, RGBColor(0xFF, 0xCC, 0x00))
add_textbox(slide, 0.5, 1.4, 5.5, 0.5,
            "  ⚠  申請進行中", font_size=16, bold=True, color=DARK_GRAY)
add_textbox(slide, 0.65, 2.05, 5.2, 1.95,
            "已與全家資訊部門確認，\n開發文件已提供。\n\n因申請流程尚未完成，\n目前無法使用全管家進行驗證。",
            font_size=14, color=DARK_GRAY)

# 狀態卡 — 完成
add_rect(slide, 6.5, 1.4, 5.5, 2.8, RGBColor(0xF0, 0xF8, 0xF2),
         line_color=GREEN)
add_rect(slide, 6.5, 1.4, 5.5, 0.5, GREEN)
add_textbox(slide, 6.5, 1.4, 5.5, 0.5,
            "  ✓  已完成項目", font_size=16, bold=True, color=WHITE)
bullet_lines(slide, 6.65, 2.05, 5.15,
             ["取得開發文件", "確認 API 規格", "確認驗證機制設計"],
             font_size=14, dot="✓  ")

# 時間軸
add_textbox(slide, 0.5, 4.55, 12.3, 0.4,
            "後續進程", font_size=16, bold=True, color=GREEN)

steps = ["申請送出", "全家審核", "測試環境開通", "正式串接", "上線驗證"]
step_w = 2.3
for i, s in enumerate(steps):
    lx = 0.5 + i * 2.52
    col = GREEN if i < 2 else RGBColor(0xBB, 0xBB, 0xBB)
    add_rect(slide, lx, 5.1, step_w, 0.55, col)
    add_textbox(slide, lx, 5.1, step_w, 0.55,
                s, font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if i < 4:
        add_textbox(slide, lx + step_w, 5.2, 0.22, 0.35,
                    "▶", font_size=14, color=MID_GRAY)

add_textbox(slide, 0.5, 5.75, 12, 0.4,
            "目前狀態：申請進行中（等待全家審核）", font_size=13,
            color=MID_GRAY, italic=True)

# ═══════════════════════════════════════════════════════════════════════════
# P5 — 角色權限對應
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout(prs))
add_rect(slide, 0, 0, 13.33, 7.5, WHITE)
slide_header(slide, "角色權限對應總覽", "Role & Permission Matrix")
slide_footer(slide, 5)

roles = ["店員", "店長", "小店長", "擔當", "管理員"]
pages = [
    ("員工首頁",       [1,1,1,0,0]),
    ("每日工作確認",   [1,1,1,0,0]),
    ("衛生自主管理",   [1,1,1,0,0]),
    ("店鋪點檢",       [1,1,1,0,0]),
    ("異常回報",       [1,1,1,0,0]),
    ("設備清潔保養",   [1,1,1,0,0]),
    ("咖啡機自檢",     [1,1,1,0,0]),
    ("月報統計",       [1,1,1,0,0]),
    ("小店長連結",     [0,1,0,0,0]),
    ("管理總覽",       [0,0,0,1,1]),
    ("店鋪狀況",       [0,0,0,1,1]),
    ("紀錄查閱",       [0,0,0,1,1]),
    ("異常管理",       [0,0,0,1,1]),
    ("數據統計",       [0,0,0,1,1]),
    ("神秘客稽查",     [0,0,0,1,1]),
]

col_widths = [2.4, 1.55, 1.55, 1.55, 1.55, 1.65]
starts_l   = [0.25, 2.72, 4.3, 5.88, 7.46, 9.04]
row_h = 0.33
header_t = 1.25

# 表頭
add_rect(slide, 0.25, header_t, 10.45, row_h, DARK_GREEN)
add_textbox(slide, starts_l[0], header_t, col_widths[0], row_h,
            "功能頁面", font_size=13, bold=True, color=WHITE)
for ci, role in enumerate(roles):
    add_textbox(slide, starts_l[ci+1], header_t, col_widths[ci+1], row_h,
                role, font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# 列
for ri, (page, perms) in enumerate(pages):
    t = header_t + (ri+1) * row_h
    bg = RGBColor(0xF0, 0xF8, 0xF2) if ri % 2 == 0 else WHITE
    add_rect(slide, 0.25, t, 10.45, row_h, bg,
             line_color=RGBColor(0xDD, 0xDD, 0xDD))
    add_textbox(slide, starts_l[0] + 0.05, t, col_widths[0], row_h,
                page, font_size=11, color=DARK_GRAY)
    for ci, perm in enumerate(perms):
        mark = "✓" if perm else "—"
        clr = GREEN if perm else RGBColor(0xBB, 0xBB, 0xBB)
        add_textbox(slide, starts_l[ci+1], t, col_widths[ci+1], row_h,
                    mark, font_size=12, bold=perm==1,
                    color=clr, align=PP_ALIGN.CENTER)

# 圖例
add_rect(slide, 10.65, 1.25, 2.35, 1.25, LIGHT_GRAY)
add_textbox(slide, 10.7, 1.28, 2.25, 0.38, "圖例", font_size=12, bold=True, color=DARK_GRAY)
add_textbox(slide, 10.7, 1.62, 2.25, 0.35,
            "✓  擁有此頁面權限", font_size=11, color=GREEN)
add_textbox(slide, 10.7, 1.95, 2.25, 0.35,
            "—  無此頁面存取權", font_size=11, color=MID_GRAY)

# ═══════════════════════════════════════════════════════════════════════════
# P6 — IOT 設備串接
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout(prs))
add_rect(slide, 0, 0, 13.33, 7.5, WHITE)
slide_header(slide, "IOT 設備串接規劃", "IoT Device Integration Plan")
slide_footer(slide, 6)

# 現況
add_rect(slide, 0.4, 1.35, 5.8, 2.65, RGBColor(0xF0, 0xF8, 0xF2),
         line_color=GREEN)
add_rect(slide, 0.4, 1.35, 5.8, 0.45, GREEN)
add_textbox(slide, 0.4, 1.35, 5.8, 0.45,
            "  全管家 IOT 現況", font_size=15, bold=True, color=WHITE)
bullet_lines(slide, 0.55, 1.9, 5.5,
             ["每天僅接收一次 IOT 設備資料",
              "資料頻率較低，無法即時反映設備狀態"],
             font_size=14, color=DARK_GRAY, line_h=0.45)

# 建議
add_rect(slide, 6.9, 1.35, 5.8, 2.65, RGBColor(0xE8, 0xF4, 0xFF),
         line_color=ACCENT)
add_rect(slide, 6.9, 1.35, 5.8, 0.45, ACCENT)
add_textbox(slide, 6.9, 1.35, 5.8, 0.45,
            "  建議整合方案", font_size=15, bold=True, color=WHITE)
bullet_lines(slide, 7.05, 1.9, 5.5,
             ["將 IOT 資料提供給每日店鋪維護系統",
              "讓店員即時確認設備狀態並填寫"],
             font_size=14, color=DARK_GRAY, line_h=0.45, dot="→  ")

# 待確認項目
add_rect(slide, 0.4, 4.25, 12.3, 2.7, RGBColor(0xFF, 0xF8, 0xE5),
         line_color=RGBColor(0xFF, 0xAA, 0x00))
add_rect(slide, 0.4, 4.25, 12.3, 0.45, RGBColor(0xFF, 0xAA, 0x00))
add_textbox(slide, 0.4, 4.25, 12.3, 0.45,
            "  ⚠  待確認議題", font_size=15, bold=True, color=WHITE)
questions = [
    "IOT 面板資料傳輸頻率（目前一天一次，是否可提高？）",
    "異常狀況觸發機制（如何即時推播或通知店員？）",
    "傳輸協議與格式規格確認（API / MQTT / Webhook？）",
]
bullet_lines(slide, 0.55, 4.82, 12.0, questions,
             font_size=14, color=DARK_GRAY, line_h=0.6, dot="❓  ")

# ═══════════════════════════════════════════════════════════════════════════
# P7 — 後台 WEB 報表
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout(prs))
add_rect(slide, 0, 0, 13.33, 7.5, WHITE)
slide_header(slide, "後台 WEB 報表", "Backend Web Reporting")
slide_footer(slide, 7)

opts = [
    ("方案 A", "資料傳輸給全家",
     "由本系統傳送資料至全家，\n由全家自行開發 WEB 報表，\n用於店鋪考核作業。",
     GREEN, "優點：整合進全家現有生態系\n考量：開發時程依賴全家排程"),
    ("方案 B", "日翊自行開發",
     "由日翊資訊獨立開發 WEB 報表，\n需符合全家對外 WEB 框架\n與權限管理規範。",
     ACCENT, "優點：靈活度高、排程自主\n考量：需符合全家對外框架規格"),
]

for i, (tag, title, desc, clr, note) in enumerate(opts):
    lx = 0.5 + i * 6.45
    add_rect(slide, lx, 1.35, 5.9, 5.55, LIGHT_GRAY,
             line_color=clr)
    add_rect(slide, lx, 1.35, 5.9, 0.6, clr)
    add_textbox(slide, lx, 1.35, 5.9, 0.6,
                f"  {tag}：{title}", font_size=16, bold=True, color=WHITE)
    add_textbox(slide, lx + 0.2, 2.1, 5.5, 2.0,
                desc, font_size=14, color=DARK_GRAY)
    add_rect(slide, lx + 0.2, 4.25, 5.5, 2.35, WHITE,
             line_color=RGBColor(0xCC, 0xCC, 0xCC))
    add_textbox(slide, lx + 0.35, 4.35, 5.2, 2.1,
                note, font_size=13, color=MID_GRAY)

add_textbox(slide, 0.5, 7.0, 12.3, 0.35,
            "* 建議本次會議確認方案選擇，以利後續開發規劃。",
            font_size=12, color=MID_GRAY, italic=True)

# ═══════════════════════════════════════════════════════════════════════════
# P8 — 資料留存時間
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout(prs))
add_rect(slide, 0, 0, 13.33, 7.5, WHITE)
slide_header(slide, "資料留存時間", "Data Retention Policy")
slide_footer(slide, 8)

rows = [
    ("每日工作確認紀錄",  "2 年",  "符合食品衛生相關法規稽查需求"),
    ("衛生自主管理紀錄",  "2 年",  "食安法規 / 衛生局稽查備查"),
    ("異常回報紀錄",      "2 年",  "問題追蹤與改善歷程留存"),
    ("設備清潔保養紀錄",  "1 年",  "設備維護管理週期對應"),
    ("咖啡機自檢紀錄",    "1 年",  "品質管理與設備保固對應"),
    ("神秘客稽查結果",    "3 年",  "總部考核依據長期保存"),
    ("IOT 設備溫度資料",  "1 年",  "設備管理 / 異常追蹤需求"),
    ("登入與存取日誌",    "6 個月","資安稽核與異常存取分析"),
]

header_cols = ["資料類型", "建議留存期限", "說明"]
col_ls = [0.35, 5.5, 7.9]
col_ws = [5.1, 2.35, 5.1]
rh = 0.52

# 表頭
add_rect(slide, 0.35, 1.3, 12.6, rh, DARK_GREEN)
for ci, h in enumerate(header_cols):
    add_textbox(slide, col_ls[ci] + 0.1, 1.3, col_ws[ci], rh,
                h, font_size=14, bold=True, color=WHITE)

for ri, (dtype, period, note) in enumerate(rows):
    t = 1.3 + (ri+1) * rh
    bg = RGBColor(0xF0, 0xF8, 0xF2) if ri % 2 == 0 else WHITE
    add_rect(slide, 0.35, t, 12.6, rh, bg,
             line_color=RGBColor(0xDD, 0xDD, 0xDD))
    add_textbox(slide, col_ls[0] + 0.1, t, col_ws[0], rh,
                dtype, font_size=12, color=DARK_GRAY)
    add_textbox(slide, col_ls[1] + 0.1, t, col_ws[1], rh,
                period, font_size=12, bold=True,
                color=GREEN, align=PP_ALIGN.CENTER)
    add_textbox(slide, col_ls[2] + 0.1, t, col_ws[2], rh,
                note, font_size=11, color=MID_GRAY)

add_textbox(slide, 0.35, 5.7, 12.6, 0.4,
            "* 具體留存期限建議與全家法遵部門確認後定案。",
            font_size=12, italic=True, color=MID_GRAY)

# ═══════════════════════════════════════════════════════════════════════════
# P9 — 下一步行動項目
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout(prs))
add_rect(slide, 0, 0, 13.33, 7.5, WHITE)
slide_header(slide, "下一步行動項目", "Next Action Items")
slide_footer(slide, 9)

actions = [
    ("01", "全管家申請跟進", "追蹤申請進度，取得測試環境帳號",
     "日翊資訊", "待全家回覆"),
    ("02", "IOT 規格確認", "確認傳輸頻率、異常通知機制與 API 格式",
     "雙方確認", "本次會議"),
    ("03", "WEB 報表方案決定", "確認由全家或日翊開發後台報表",
     "雙方決策", "本次會議"),
    ("04", "資料留存政策確認", "與全家法遵確認各類資料留存期限",
     "全家法遵", "下次會議前"),
    ("05", "前台改善項目開發", "完成五項改善功能的開發與測試",
     "日翊資訊", "2026/05"),
]

col_ls2  = [0.35, 1.15, 3.5, 8.7, 10.8]
col_ws2  = [0.75, 2.3, 5.15, 2.05, 2.1]
headers2 = ["#", "項目", "說明", "負責方", "目標時程"]
rh2 = 0.52

add_rect(slide, 0.35, 1.3, 12.6, rh2, DARK_GREEN)
for ci, h in enumerate(headers2):
    add_textbox(slide, col_ls2[ci] + 0.05, 1.3, col_ws2[ci], rh2,
                h, font_size=14, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT)

for ri, (no, title, desc, owner, deadline) in enumerate(actions):
    t = 1.3 + (ri+1) * rh2
    bg = RGBColor(0xF0, 0xF8, 0xF2) if ri % 2 == 0 else WHITE
    add_rect(slide, 0.35, t, 12.6, rh2, bg,
             line_color=RGBColor(0xDD, 0xDD, 0xDD))
    add_rect(slide, col_ls2[0], t, col_ws2[0], rh2, GREEN)
    add_textbox(slide, col_ls2[0], t, col_ws2[0], rh2,
                no, font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, col_ls2[1] + 0.05, t, col_ws2[1], rh2,
                title, font_size=13, bold=True, color=DARK_GREEN)
    add_textbox(slide, col_ls2[2] + 0.05, t, col_ws2[2], rh2,
                desc, font_size=12, color=DARK_GRAY)
    add_textbox(slide, col_ls2[3], t, col_ws2[3], rh2,
                owner, font_size=12, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, col_ls2[4], t, col_ws2[4], rh2,
                deadline, font_size=12, color=MID_GRAY, align=PP_ALIGN.CENTER)

# ── 儲存 ────────────────────────────────────────────────────────────────────
output_path = r"D:\claude_cli\store-management-app\docs\每日店鋪維護系統_開發討論會議_20260415.pptx"
prs.save(output_path)
print(f"已儲存：{output_path}")
